"""
Generate Max-Readability PPTX Slide for Nyaya Setu: Tech Stack & System Architecture
- Extra large fonts for maximum projector visibility
- Harmonious light color palette on right side (no dark blue)
- Simplified text & details, skipped port numbers
- Exact header match: | TECH STACK & SYSTEM ARCHITECTURE
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_max_readability_pptx():
    prs = Presentation()
    # 16:9 widescreen dimensions (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 1. Slide Title Header (Exact match: Orange vertical bar + TECH STACK & SYSTEM ARCHITECTURE)
    orange_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.55), Inches(0.4), Inches(0.14), Inches(0.72)
    )
    orange_bar.fill.solid()
    orange_bar.fill.fore_color.rgb = RGBColor(249, 115, 22) # #F97316
    orange_bar.line.fill.background()

    # Title Text
    title_box = slide.shapes.add_textbox(Inches(0.85), Inches(0.32), Inches(11.0), Inches(0.85))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "TECH STACK & SYSTEM ARCHITECTURE"
    p.font.name = "Segoe UI"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)

    # 2. LEFT CARD: Tech Stack Breakdown
    left_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.55), Inches(1.2), Inches(5.95), Inches(5.65)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    left_card.line.color.rgb = RGBColor(226, 232, 240)
    left_card.line.width = Pt(2)

    # Left Header: LAYER & PRODUCTION TECHNOLOGY
    lh_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(5.5), Inches(0.4))
    lh_tf = lh_box.text_frame
    lh_p = lh_tf.paragraphs[0]
    lh_p.text = "● LAYER                                         PRODUCTION TECHNOLOGY"
    lh_p.font.name = "Segoe UI"
    lh_p.font.size = Pt(12)
    lh_p.font.bold = True
    lh_p.font.color.rgb = RGBColor(71, 85, 105)

    # Left Rows
    tech_items = [
        ("Frontend & UI", "React 18 • Tailwind CSS • Vite", RGBColor(238, 242, 255), RGBColor(67, 56, 202)),
        ("App Gateway", "Node.js • Express • RBAC • JWT", RGBColor(236, 253, 245), RGBColor(4, 120, 87)),
        ("AI Microservice", "Python 3.11 • FastAPI • Uvicorn", RGBColor(236, 254, 255), RGBColor(14, 116, 144)),
        ("Agent Engine", "LangGraph • 11 Specialized Agents", RGBColor(250, 245, 255), RGBColor(126, 34, 206)),
        ("Guardrail & Trust", "8-Tier Active Security Perimeter", RGBColor(254, 252, 232), RGBColor(161, 98, 7)),
        ("Hybrid Legal RAG", "Dense Vectors + BM25 + Reranker", RGBColor(239, 246, 255), RGBColor(29, 78, 216)),
        ("DB & Caching", "MongoDB 7.0 • Redis 7.2 Queues", RGBColor(255, 241, 242), RGBColor(190, 18, 60)),
        ("Doc AI & Voice", "Clause NLP • Whisper STT • jsPDF", RGBColor(240, 253, 250), RGBColor(15, 118, 110)),
        ("Deployment", "Docker Compose • 5 Isolated Services", RGBColor(241, 245, 249), RGBColor(51, 65, 85)),
    ]

    row_y = 1.8
    for label, tech, pill_bg, pill_fg in tech_items:
        # Label
        lbl_box = slide.shapes.add_textbox(Inches(0.75), Inches(row_y), Inches(2.4), Inches(0.4))
        lbl_tf = lbl_box.text_frame
        lbl_p = lbl_tf.paragraphs[0]
        lbl_p.text = label
        lbl_p.font.name = "Segoe UI"
        lbl_p.font.size = Pt(14)
        lbl_p.font.bold = True
        lbl_p.font.color.rgb = RGBColor(15, 23, 42)

        # Pill
        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(3.2), Inches(row_y + 0.02), Inches(3.15), Inches(0.4)
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = pill_bg
        pill.line.fill.background()
        
        pill_tf = pill.text_frame
        pill_tf.word_wrap = False
        pill_p = pill_tf.paragraphs[0]
        pill_p.text = tech
        pill_p.alignment = PP_ALIGN.CENTER
        pill_p.font.name = "Segoe UI"
        pill_p.font.size = Pt(12)
        pill_p.font.bold = True
        pill_p.font.color.rgb = pill_fg

        row_y += 0.55

    # 3. RIGHT CARD (Revamped Harmonious Light Flow)
    right_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.2), Inches(5.95), Inches(5.65)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    right_card.line.color.rgb = RGBColor(226, 232, 240)
    right_card.line.width = Pt(2)

    # Right Header
    rh_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.4))
    rh_tf = rh_box.text_frame
    rh_p = rh_tf.paragraphs[0]
    rh_p.text = "● SYSTEM ARCHITECTURE & DATA FLOW"
    rh_p.font.name = "Segoe UI"
    rh_p.font.size = Pt(12)
    rh_p.font.bold = True
    rh_p.font.color.rgb = RGBColor(71, 85, 105)

    # Tier 1 (Light Blue)
    t1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.72), Inches(5.55), Inches(0.75))
    t1.fill.solid()
    t1.fill.fore_color.rgb = RGBColor(240, 249, 255)
    t1.line.color.rgb = RGBColor(186, 230, 253)
    t1.line.width = Pt(1.5)
    t1_tf = t1.text_frame
    t1_p = t1_tf.paragraphs[0]
    t1_p.text = "1. React 18 + Tailwind Client Portal"
    t1_p.font.size = Pt(13)
    t1_p.font.bold = True
    t1_p.font.color.rgb = RGBColor(12, 74, 110)
    t1_sub = t1_tf.add_paragraph()
    t1_sub.text = "Citizen Intake UI • Lawyer Workspace • Live Trust Badges • Voice Assistant"
    t1_sub.font.size = Pt(10)
    t1_sub.font.color.rgb = RGBColor(71, 85, 105)

    # Arrow 1
    a1 = slide.shapes.add_textbox(Inches(7.0), Inches(2.48), Inches(5.55), Inches(0.25))
    a1_tf = a1.text_frame
    a1_p = a1_tf.paragraphs[0]
    a1_p.text = "↓ REST API / HTTPS + JWT Authentication"
    a1_p.alignment = PP_ALIGN.CENTER
    a1_p.font.size = Pt(9.5)
    a1_p.font.bold = True
    a1_p.font.color.rgb = RGBColor(148, 163, 184)

    # Tier 2 (Light Emerald)
    t2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(2.72), Inches(5.55), Inches(0.75))
    t2.fill.solid()
    t2.fill.fore_color.rgb = RGBColor(236, 253, 245)
    t2.line.color.rgb = RGBColor(167, 243, 208)
    t2.line.width = Pt(1.5)
    t2_tf = t2.text_frame
    t2_p = t2_tf.paragraphs[0]
    t2_p.text = "2. Node.js + Express Application Gateway"
    t2_p.font.size = Pt(13)
    t2_p.font.bold = True
    t2_p.font.color.rgb = RGBColor(6, 78, 59)
    t2_sub = t2_tf.add_paragraph()
    t2_sub.text = "Granular RBAC • Tenant Case Isolation • Redis Task Queues • Audit Logger"
    t2_sub.font.size = Pt(10)
    t2_sub.font.color.rgb = RGBColor(71, 85, 105)

    # Arrow 2
    a2 = slide.shapes.add_textbox(Inches(7.0), Inches(3.48), Inches(5.55), Inches(0.25))
    a2_tf = a2.text_frame
    a2_p = a2_tf.paragraphs[0]
    a2_p.text = "↓ Microservice Proxy (FastAPI Core)"
    a2_p.alignment = PP_ALIGN.CENTER
    a2_p.font.size = Pt(9.5)
    a2_p.font.bold = True
    a2_p.font.color.rgb = RGBColor(148, 163, 184)

    # Tier 3 (Light Purple Container)
    t3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(3.72), Inches(5.55), Inches(1.8))
    t3.fill.solid()
    t3.fill.fore_color.rgb = RGBColor(250, 245, 255)
    t3.line.color.rgb = RGBColor(221, 214, 254)
    t3.line.width = Pt(1.5)

    # T3 Title
    t3_title = slide.shapes.add_textbox(Inches(7.1), Inches(3.77), Inches(5.35), Inches(0.32))
    t3_tf = t3_title.text_frame
    t3_p = t3_tf.paragraphs[0]
    t3_p.text = "3. Python FastAPI + LangGraph Multi-Agent Engine"
    t3_p.font.size = Pt(13)
    t3_p.font.bold = True
    t3_p.font.color.rgb = RGBColor(88, 28, 135)

    # Agent Badges
    agent_badges = [
        ("Intake", Inches(7.15), Inches(4.14), Inches(0.83)),
        ("Research", Inches(8.04), Inches(4.14), Inches(0.86)),
        ("Doc AI", Inches(8.96), Inches(4.14), Inches(0.83)),
        ("Risk 1930", Inches(9.85), Inches(4.14), Inches(0.88)),
        ("Drafting", Inches(10.79), Inches(4.14), Inches(0.83)),
        ("Verifier", Inches(11.68), Inches(4.14), Inches(0.83)),
    ]

    for title, x, y, w in agent_badges:
        ab = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.33))
        ab.fill.solid()
        ab.fill.fore_color.rgb = RGBColor(255, 255, 255)
        ab.line.color.rgb = RGBColor(203, 213, 225)
        ab.line.width = Pt(1)
        ab_tf = ab.text_frame
        ab_p = ab_tf.paragraphs[0]
        ab_p.text = title
        ab_p.alignment = PP_ALIGN.CENTER
        ab_p.font.size = Pt(10)
        ab_p.font.bold = True
        ab_p.font.color.rgb = RGBColor(30, 41, 59)

    # Guardrail Shield Banner (Light Amber)
    gb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.15), Inches(4.56), Inches(5.25), Inches(0.72))
    gb.fill.solid()
    gb.fill.fore_color.rgb = RGBColor(254, 252, 232)
    gb.line.color.rgb = RGBColor(253, 224, 71)
    gb.line.width = Pt(1.5)
    gb_tf = gb.text_frame
    gb_p = gb_tf.paragraphs[0]
    gb_p.text = "🛡️ Active Guardrail Perimeter (Zero-Trust AI Safety Layer)"
    gb_p.font.size = Pt(11)
    gb_p.font.bold = True
    gb_p.font.color.rgb = RGBColor(161, 98, 7)
    gb_sub = gb_tf.add_paragraph()
    gb_sub.text = "PII Redaction (Aadhaar/PAN) • Prompt Shield • 1930 Emergency Helpline"
    gb_sub.font.size = Pt(9.5)
    gb_sub.font.color.rgb = RGBColor(133, 77, 14)

    # Arrow 3
    a3 = slide.shapes.add_textbox(Inches(7.0), Inches(5.53), Inches(5.55), Inches(0.25))
    a3_tf = a3.text_frame
    a3_p = a3_tf.paragraphs[0]
    a3_p.text = "↓ Dense Vectors + BM25 Hybrid Retrieval"
    a3_p.alignment = PP_ALIGN.CENTER
    a3_p.font.size = Pt(9.5)
    a3_p.font.bold = True
    a3_p.font.color.rgb = RGBColor(148, 163, 184)

    # Tier 4 (Light Amber)
    t4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(5.76), Inches(5.55), Inches(0.75))
    t4.fill.solid()
    t4.fill.fore_color.rgb = RGBColor(255, 251, 235)
    t4.line.color.rgb = RGBColor(253, 224, 71)
    t4.line.width = Pt(1.5)
    t4_tf = t4.text_frame
    t4_p = t4_tf.paragraphs[0]
    t4_p.text = "4. Tier-1 Statutory Knowledge Base & Gazette Roll"
    t4_p.font.size = Pt(13)
    t4_p.font.bold = True
    t4_p.font.color.rgb = RGBColor(120, 53, 15)
    t4_sub = t4_tf.add_paragraph()
    t4_sub.text = "India Code • Gazette Notifications • eCourts Records • NALSA Legal Aid Portal"
    t4_sub.font.size = Pt(10)
    t4_sub.font.color.rgb = RGBColor(146, 64, 14)

    # 4. Bottom Footer Bar
    footer_box = slide.shapes.add_textbox(Inches(0.55), Inches(6.94), Inches(12.2), Inches(0.4))
    ft_tf = footer_box.text_frame
    ft_p = ft_tf.paragraphs[0]
    ft_p.text = "BUILD WITH भारत 2.0  •  NATIONAL LEVEL HACKATHON                  TEAM CODING CHAMPS • BHARATI VIDYAPEETH'S COE"
    ft_p.font.name = "Segoe UI"
    ft_p.font.size = Pt(11)
    ft_p.font.bold = True
    ft_p.font.color.rgb = RGBColor(100, 116, 139)

    out_pptx = "Tech_Stack_Slide_Revamped.pptx"
    prs.save(out_pptx)
    print(f"Max Readability PPTX saved to {out_pptx}")

if __name__ == "__main__":
    create_max_readability_pptx()
