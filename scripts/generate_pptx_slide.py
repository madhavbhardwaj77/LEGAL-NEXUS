"""
Generate High-Quality PPTX Slide for Nyaya Setu: Tech Stack & System Architecture
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_tech_stack_presentation():
    prs = Presentation()
    # Set 16:9 widescreen dimensions (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # 1. Slide Title Header (Exact match: Orange vertical bar + TECH STACK & SYSTEM ARCHITECTURE)
    # Orange vertical bar
    orange_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(0.5), Inches(0.12), Inches(0.65)
    )
    orange_bar.fill.solid()
    orange_bar.fill.fore_color.rgb = RGBColor(249, 115, 22) # Saffron Orange #F97316
    orange_bar.line.fill.background()

    # Title Text
    title_box = slide.shapes.add_textbox(Inches(1.05), Inches(0.42), Inches(9.5), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "TECH STACK & SYSTEM ARCHITECTURE"
    p.font.name = "Arial"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42) # Slate 900

    # 2. LEFT CARD: Tech Stack Breakdown (Approved Part)
    left_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.35), Inches(5.6), Inches(5.4)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    left_card.line.color.rgb = RGBColor(226, 232, 240)
    left_card.line.width = Pt(1.5)

    # Left Card Header: LAYER & PRODUCTION TECHNOLOGY
    lh_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.45), Inches(5.2), Inches(0.4))
    lh_tf = lh_box.text_frame
    lh_p = lh_tf.paragraphs[0]
    lh_p.text = "● LAYER                                    PRODUCTION TECHNOLOGY"
    lh_p.font.name = "Arial"
    lh_p.font.size = Pt(10)
    lh_p.font.bold = True
    lh_p.font.color.rgb = RGBColor(71, 85, 105)

    # Left Card Rows
    tech_items = [
        ("Frontend & UI", "React 18 • Tailwind CSS • Vite", RGBColor(238, 242, 255), RGBColor(79, 70, 229)),
        ("App Gateway", "Node.js • Express • RBAC • JWT", RGBColor(236, 253, 245), RGBColor(5, 150, 105)),
        ("AI Microservice", "Python 3.11 • FastAPI • Uvicorn", RGBColor(236, 254, 255), RGBColor(8, 145, 178)),
        ("Agent Engine", "LangGraph • 11 Specialized Agents", RGBColor(250, 245, 255), RGBColor(147, 51, 234)),
        ("Guardrail & Trust", "8-Tier Active Security Perimeter", RGBColor(254, 252, 232), RGBColor(202, 138, 4)),
        ("Hybrid Legal RAG", "Dense Vectors + BM25 + Reranker", RGBColor(239, 246, 255), RGBColor(37, 99, 235)),
        ("DB & Caching", "MongoDB 7.0 • Redis 7.2 Queues", RGBColor(255, 241, 242), RGBColor(225, 29, 72)),
        ("Doc AI & Voice", "Clause NLP • Whisper STT • jsPDF", RGBColor(240, 253, 250), RGBColor(13, 148, 136)),
        ("Deployment", "Docker Compose • 5 Isolated Services", RGBColor(241, 245, 249), RGBColor(71, 85, 105)),
    ]

    row_y = 1.95
    for label, tech, pill_bg, pill_text_col in tech_items:
        # Label
        lbl_box = slide.shapes.add_textbox(Inches(1.0), Inches(row_y), Inches(2.2), Inches(0.4))
        lbl_tf = lbl_box.text_frame
        lbl_p = lbl_tf.paragraphs[0]
        lbl_p.text = label
        lbl_p.font.name = "Arial"
        lbl_p.font.size = Pt(11)
        lbl_p.font.bold = True
        lbl_p.font.color.rgb = RGBColor(15, 23, 42)

        # Pill background
        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(3.1), Inches(row_y + 0.04), Inches(3.1), Inches(0.32)
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = pill_bg
        pill.line.fill.background()
        
        pill_tf = pill.text_frame
        pill_tf.word_wrap = False
        pill_p = pill_tf.paragraphs[0]
        pill_p.text = tech
        pill_p.alignment = PP_ALIGN.CENTER
        pill_p.font.name = "Arial"
        pill_p.font.size = Pt(9.5)
        pill_p.font.bold = True
        pill_p.font.color.rgb = pill_text_col

        row_y += 0.52

    # 3. RIGHT CARD: Redesigned Architecture Workflow
    right_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.35), Inches(5.7), Inches(5.4)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    right_card.line.color.rgb = RGBColor(226, 232, 240)
    right_card.line.width = Pt(1.5)

    # Right Card Header
    rh_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.45), Inches(5.3), Inches(0.4))
    rh_tf = rh_box.text_frame
    rh_p = rh_tf.paragraphs[0]
    rh_p.text = "● SYSTEM ARCHITECTURE & DATA PIPELINE"
    rh_p.font.name = "Arial"
    rh_p.font.size = Pt(10)
    rh_p.font.bold = True
    rh_p.font.color.rgb = RGBColor(71, 85, 105)

    # Tier 1: Client Layer Box
    t1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.0), Inches(1.85), Inches(5.3), Inches(0.65)
    )
    t1.fill.solid()
    t1.fill.fore_color.rgb = RGBColor(248, 250, 252)
    t1.line.color.rgb = RGBColor(191, 219, 254) # Blue border
    t1.line.width = Pt(1.2)
    t1_tf = t1.text_frame
    t1_p = t1_tf.paragraphs[0]
    t1_p.text = "1. Client Layer: React 18 + Tailwind Responsive App (Port 5173)"
    t1_p.font.size = Pt(10.5)
    t1_p.font.bold = True
    t1_p.font.color.rgb = RGBColor(30, 58, 138)
    t1_sub = t1_tf.add_paragraph()
    t1_sub.text = "Citizen Story Intake • Lawyer Workspace • Live Trust Badge UX • Voice Widget"
    t1_sub.font.size = Pt(8.5)
    t1_sub.font.color.rgb = RGBColor(100, 116, 139)

    # Arrow 1
    a1 = slide.shapes.add_textbox(Inches(7.0), Inches(2.5), Inches(5.3), Inches(0.25))
    a1_tf = a1.text_frame
    a1_p = a1_tf.paragraphs[0]
    a1_p.text = "↓ REST API / HTTPS + JWT Authentication"
    a1_p.alignment = PP_ALIGN.CENTER
    a1_p.font.size = Pt(8.5)
    a1_p.font.bold = True
    a1_p.font.color.rgb = RGBColor(148, 163, 184)

    # Tier 2: Gateway Layer Box
    t2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.0), Inches(2.75), Inches(5.3), Inches(0.65)
    )
    t2.fill.solid()
    t2.fill.fore_color.rgb = RGBColor(248, 250, 252)
    t2.line.color.rgb = RGBColor(167, 243, 208) # Emerald border
    t2.line.width = Pt(1.2)
    t2_tf = t2.text_frame
    t2_p = t2_tf.paragraphs[0]
    t2_p.text = "2. Gateway Layer: Node.js + Express REST Core (Port 5001)"
    t2_p.font.size = Pt(10.5)
    t2_p.font.bold = True
    t2_p.font.color.rgb = RGBColor(6, 95, 70)
    t2_sub = t2_tf.add_paragraph()
    t2_sub.text = "Granular RBAC • Tenant Isolation • Redis Task Queue • Central Audit Logger"
    t2_sub.font.size = Pt(8.5)
    t2_sub.font.color.rgb = RGBColor(100, 116, 139)

    # Arrow 2
    a2 = slide.shapes.add_textbox(Inches(7.0), Inches(3.4), Inches(5.3), Inches(0.25))
    a2_tf = a2.text_frame
    a2_p = a2_tf.paragraphs[0]
    a2_p.text = "↓ Internal Microservice Proxy (FastAPI Router)"
    a2_p.alignment = PP_ALIGN.CENTER
    a2_p.font.size = Pt(8.5)
    a2_p.font.bold = True
    a2_p.font.color.rgb = RGBColor(148, 163, 184)

    # Tier 3: AI Engine & Multi-Agent Core (Elevated Container)
    t3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.0), Inches(3.65), Inches(5.3), Inches(1.7)
    )
    t3.fill.solid()
    t3.fill.fore_color.rgb = RGBColor(15, 23, 42) # Slate 900 dark container
    t3.line.color.rgb = RGBColor(51, 65, 85)
    t3.line.width = Pt(1.2)

    # T3 Title
    t3_title = slide.shapes.add_textbox(Inches(7.1), Inches(3.7), Inches(5.1), Inches(0.3))
    t3_tf = t3_title.text_frame
    t3_p = t3_tf.paragraphs[0]
    t3_p.text = "3. AI Core: Python FastAPI + LangGraph Multi-Agent Engine (Port 8001)"
    t3_p.font.size = Pt(9.5)
    t3_p.font.bold = True
    t3_p.font.color.rgb = RGBColor(226, 232, 240)

    # Specialized Agent Pills Grid
    agent_badges = [
        ("🗣️ Intake", Inches(7.15), Inches(4.05), Inches(1.55)),
        ("⚖️ Research", Inches(2.7 + 5.3*0.35 + 0.1), Inches(4.05), Inches(1.6)),
        ("📄 Doc AI", Inches(2.7 + 5.3*0.35 + 1.8), Inches(4.05), Inches(1.55)),
        ("🚨 Risk 1930", Inches(7.15), Inches(4.45), Inches(1.55)),
        ("✍️ Drafting", Inches(2.7 + 5.3*0.35 + 0.1), Inches(4.45), Inches(1.6)),
        ("✅ Verifier", Inches(2.7 + 5.3*0.35 + 1.8), Inches(4.45), Inches(1.55)),
    ]

    for title, x, y, w in agent_badges:
        ab = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.32))
        ab.fill.solid()
        ab.fill.fore_color.rgb = RGBColor(30, 41, 59) # Slate 800
        ab.line.color.rgb = RGBColor(71, 85, 105)
        ab_tf = ab.text_frame
        ab_p = ab_tf.paragraphs[0]
        ab_p.text = title
        ab_p.alignment = PP_ALIGN.CENTER
        ab_p.font.size = Pt(8.5)
        ab_p.font.bold = True
        ab_p.font.color.rgb = RGBColor(241, 245, 249)

    # Guardrail Banner
    gb = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.15), Inches(4.88), Inches(5.0), Inches(0.35)
    )
    gb.fill.solid()
    gb.fill.fore_color.rgb = RGBColor(69, 26, 3) # Amber dark bg
    gb.line.color.rgb = RGBColor(217, 119, 6) # Amber border
    gb_tf = gb.text_frame
    gb_p = gb_tf.paragraphs[0]
    gb_p.text = "🛡️ Active Guardrail Perimeter: PII Redaction • Prompt Shield • 1930 Routing"
    gb_p.alignment = PP_ALIGN.CENTER
    gb_p.font.size = Pt(8.5)
    gb_p.font.bold = True
    gb_p.font.color.rgb = RGBColor(254, 240, 138)

    # Arrow 3
    a3 = slide.shapes.add_textbox(Inches(7.0), Inches(5.35), Inches(5.3), Inches(0.25))
    a3_tf = a3.text_frame
    a3_p = a3_tf.paragraphs[0]
    a3_p.text = "↓ Hybrid Retrieval (Dense Vectors + BM25 Lexical + Cross-Encoder)"
    a3_p.alignment = PP_ALIGN.CENTER
    a3_p.font.size = Pt(8.5)
    a3_p.font.bold = True
    a3_p.font.color.rgb = RGBColor(148, 163, 184)

    # Tier 4: Knowledge Base Grounding Box
    t4 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.0), Inches(5.6), Inches(5.3), Inches(0.65)
    )
    t4.fill.solid()
    t4.fill.fore_color.rgb = RGBColor(254, 252, 232)
    t4.line.color.rgb = RGBColor(253, 224, 71) # Gold/amber border
    t4.line.width = Pt(1.2)
    t4_tf = t4.text_frame
    t4_p = t4_tf.paragraphs[0]
    t4_p.text = "4. Grounding Base: Tier-1 Official Gazette & Statutory Rolls"
    t4_p.font.size = Pt(10.5)
    t4_p.font.bold = True
    t4_p.font.color.rgb = RGBColor(133, 77, 14)
    t4_sub = t4_tf.add_paragraph()
    t4_sub.text = "India Code • e-Gazette Notifications • eCourts Records • NALSA Legal Aid Portal"
    t4_sub.font.size = Pt(8.5)
    t4_sub.font.color.rgb = RGBColor(161, 98, 7)

    # 4. Bottom Footer Bar (Matching Hackathon Template)
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.4))
    ft_tf = footer_box.text_frame
    ft_p = ft_tf.paragraphs[0]
    ft_p.text = "BUILD WITH भारत 2.0                                             TEAM CODING CHAMPS • BHARATI VIDYAPEETH'S COLLEGE OF ENGINEERING"
    ft_p.font.name = "Arial"
    ft_p.font.size = Pt(9.5)
    ft_p.font.bold = True
    ft_p.font.color.rgb = RGBColor(100, 116, 139)

    # Save presentation
    output_path = "Tech_Stack_Slide.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_tech_stack_presentation()
